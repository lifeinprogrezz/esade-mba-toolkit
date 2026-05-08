"""
ESADE deck design system + reusable visual helpers.

Import everything you need from this module when building a case deck:

    from esade_design import *

Available primitives are listed by section below. See velos_deck.py for a
complete worked example.

DESIGN PHILOSOPHY
-----------------
- 16:9 widescreen, navy + grey + one accent (deep red).
- One typeface (Calibri) — no decoration, no clipart, no AI-generated images.
- Each helper draws exactly one element; compose them to build a slide.
- The professor will spot generic AI aesthetics in five seconds.
  Always frame slides in the course's vocabulary, not generic strategy-speak.
"""

from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib.colors import HexColor
from reportlab.platypus import SimpleDocTemplate
from reportlab.lib.enums import TA_LEFT, TA_JUSTIFY


# =============================================================================
# COLOR PALETTE
# =============================================================================

NAVY      = RGBColor(0x0A, 0x25, 0x40)   # primary brand
NAVY_LITE = RGBColor(0x1B, 0x3A, 0x57)
ACCENT    = RGBColor(0xC0, 0x39, 0x2B)   # deep red — call-outs, warnings
AMBER     = RGBColor(0xD9, 0x77, 0x06)
GREEN     = RGBColor(0x05, 0x96, 0x69)
RED       = RGBColor(0xDC, 0x26, 0x26)
GREY_DK   = RGBColor(0x1F, 0x29, 0x37)
GREY      = RGBColor(0x6B, 0x72, 0x80)
GREY_LT   = RGBColor(0xE5, 0xE7, 0xEB)
BG_LIGHT  = RGBColor(0xF9, 0xFA, 0xFB)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)


# =============================================================================
# TYPE & DIMENSIONS
# =============================================================================

FONT = "Calibri"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# =============================================================================
# CORE BUILDERS
# =============================================================================

def setup_presentation():
    """Return an empty 16:9 Presentation, ready for blank slides."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs):
    """Add a fully blank slide (layout 6) — no placeholders."""
    return prs.slides.add_slide(prs.slide_layouts[6])


# =============================================================================
# SHAPE PRIMITIVES
# =============================================================================

def fill_shape(shape, rgb):
    """Solid fill, no border."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, fill=None, line=None):
    """Add a rectangle. fill=None means transparent. line=None means no border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text,
             font=FONT, size=14, bold=False, color=GREY_DK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    """Add a textbox. `text` may be a string or a list of strings (one per line)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    if isinstance(text, str):
        text = [text]

    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_rich_text(slide, left, top, width, height, runs,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """Single-paragraph textbox with mixed formatting.

    runs = [(text, {size, bold, color, italic, font}), ...]
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for txt, props in runs:
        run = p.add_run()
        run.text = txt
        run.font.name = props.get("font", FONT)
        run.font.size = Pt(props.get("size", 14))
        run.font.bold = props.get("bold", False)
        run.font.italic = props.get("italic", False)
        run.font.color.rgb = props.get("color", GREY_DK)
    return tb


def add_bullets(slide, left, top, width, height, items,
                size=13, color=GREY_DK, bullet_color=NAVY,
                line_spacing=1.25, bold_first_word=False):
    """Bullet list. Each item is either a string or a (bold_part, rest) tuple."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run_b = p.add_run()
        run_b.text = "•  "
        run_b.font.name = FONT
        run_b.font.size = Pt(size)
        run_b.font.bold = True
        run_b.font.color.rgb = bullet_color
        if isinstance(item, tuple):
            bold_part, rest = item
            r1 = p.add_run()
            r1.text = bold_part
            r1.font.name = FONT
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = rest
            r2.font.name = FONT
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


# =============================================================================
# SLIDE CHROME (TITLE BAR + FOOTER)
# =============================================================================

def slide_chrome(slide, title, subtitle=None, page_num=None, total=None):
    """Standard navy title bar + grey subtitle line + slide number footer."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.85), fill=NAVY)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(11.5), Inches(0.5),
             title, size=22, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4),
                 subtitle, size=13, italic=True, color=GREY)
    if page_num is not None:
        add_text(slide, Inches(12.0), Inches(7.15), Inches(1.2), Inches(0.3),
                 f"{page_num}/{total}",
                 size=9, color=GREY, align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.5), Inches(7.10), Inches(12.3), Emu(8000),
             fill=GREY_LT)


def add_callout(slide, left, top, width, height, text,
                fill=NAVY, text_color=WHITE, size=13, italic=True, bold=False):
    """Solid-color band for a one-liner takeaway."""
    add_rect(slide, left, top, width, height, fill=fill)
    add_text(slide, left + Inches(0.2), top + Inches(0.05),
             width - Inches(0.4), height - Inches(0.1),
             text, size=size, italic=italic, bold=bold,
             color=text_color, anchor=MSO_ANCHOR.MIDDLE)


# =============================================================================
# TABLES
# =============================================================================

def add_table(slide, left, top, width, height, data,
              header_fill=NAVY, header_color=WHITE,
              row_fill=BG_LIGHT, alt_fill=WHITE,
              first_col_bold=False, first_col_fill=None,
              size=11, header_size=12, col_widths=None):
    """Comparison/data table. `data` is a list of rows; row 0 is the header."""
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                font_color = header_color
                fsize = header_size
                fbold = True
            else:
                if first_col_bold and c == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = first_col_fill or NAVY_LITE
                    font_color = WHITE
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_fill if r % 2 == 0 else row_fill
                    font_color = GREY_DK
                fsize = size
                fbold = (first_col_bold and c == 0)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = ""
            run = p.add_run()
            run.text = data[r][c]
            run.font.name = FONT
            run.font.size = Pt(fsize)
            run.font.bold = fbold
            run.font.color.rgb = font_color
    return table_shape


# =============================================================================
# AUTO-FIT GUARD
# =============================================================================

def _disable_autosize(shape):
    """Force shape NOT to auto-shrink text or auto-resize to text.

    Use this before adding text to non-rectangular shapes (chevrons, pentagons,
    arrows) — their built-in text frames otherwise auto-shrink and produce
    vertical character wrapping.
    """
    sp = shape.text_frame._txBody
    bodyPr = sp.find(qn("a:bodyPr"))
    if bodyPr is not None:
        for child in list(bodyPr):
            tag = etree.QName(child.tag).localname
            if tag in ("normAutofit", "spAutoFit", "noAutofit"):
                bodyPr.remove(child)
        bodyPr.append(etree.SubElement(bodyPr, qn("a:noAutofit")))


# =============================================================================
# DOMAIN-SPECIFIC VISUAL PRIMITIVES
# =============================================================================
# These are tuned for the Tech & Digital Business course's vocabulary.
# Other courses may need their own primitives — add them here as you go.

def draw_inevitability_cascade(slide, left, top, width, height,
                               statuses, labels, sublabels=None,
                               highlight_last=False, compact=False):
    """4-phase Inevitability Cascade (Tech & Digital Business framework).

    statuses: list of 4 in {"ok", "warn", "fail"} — drives color + icon.
    labels:   list of 4 phase names.
    sublabels: optional list of 4 short captions.
    highlight_last: force the 4th block to ACCENT regardless of status.
    compact: use plain rectangles + arrow separators instead of pentagons.
             Required for tight (<2.0in tall, <8in wide) layouts.
    """
    n = 4
    icon_map = {"ok": "✓", "warn": "!", "fail": "✗"}

    if compact:
        sep_w = Inches(0.18)
        gap = Inches(0.05)
        total_sep = (sep_w + gap * 2) * (n - 1)
        block_w = (width - total_sep) / n
        margin_h = Inches(0.05)
        margin_v = Inches(0.05)
        icon_size = 11
        label_size = 11
    else:
        gap = Inches(0.1)
        block_w = (width - gap * (n - 1)) / n
        margin_h_l = Inches(0.15)
        margin_h_r = Inches(0.55)
        margin_v = Inches(0.12)
        icon_size = 16
        label_size = 13

    cursor_x = left
    for i in range(n):
        if compact:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, cursor_x, top, block_w, height
            )
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.PENTAGON, cursor_x, top, block_w, height
            )

        if i == n - 1 and highlight_last:
            fill_shape(shape, ACCENT)
        else:
            status = statuses[i]
            if status == "ok":
                fill_shape(shape, NAVY)
            elif status == "warn":
                fill_shape(shape, AMBER)
            elif status == "fail":
                fill_shape(shape, ACCENT)
            else:
                fill_shape(shape, NAVY_LITE)

        tf = shape.text_frame
        tf.word_wrap = True
        if compact:
            tf.margin_left = margin_h
            tf.margin_right = margin_h
        else:
            tf.margin_left = margin_h_l
            tf.margin_right = margin_h_r
        tf.margin_top = margin_v
        tf.margin_bottom = margin_v
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _disable_autosize(shape)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if compact else PP_ALIGN.LEFT

        ico = icon_map.get(statuses[i], "•")
        r1 = p.add_run()
        r1.text = f"{ico}  "
        r1.font.name = FONT
        r1.font.size = Pt(icon_size)
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        r2 = p.add_run()
        r2.text = labels[i]
        r2.font.name = FONT
        r2.font.size = Pt(label_size)
        r2.font.bold = True
        r2.font.color.rgb = WHITE

        if sublabels:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER if compact else PP_ALIGN.LEFT
            r = p2.add_run()
            r.text = sublabels[i]
            r.font.name = FONT
            r.font.size = Pt(9)
            r.font.color.rgb = WHITE

        cursor_x = cursor_x + block_w
        if compact and i < n - 1:
            arr_w = Inches(0.18)
            arr_h = Inches(0.22)
            arr_x = cursor_x + gap
            arr_y = top + (height - arr_h) / 2
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arr_x, arr_y, arr_w, arr_h
            )
            fill_shape(arr, GREY)
            cursor_x = arr_x + arr_w + gap


def draw_layer_arrows(slide, left, top, width, height, layers):
    """Horizontal progress arrows for showing kinetic asymmetries between layers.

    layers = [(name, progress 0..1, color), ...]
    Lightning-bolt marker appears next to layers with progress < 0.5.
    """
    n = len(layers)
    row_h = height / n
    for i, (name, progress, color) in enumerate(layers):
        y = top + row_h * i + Inches(0.1)
        h = row_h - Inches(0.2)
        add_text(slide, left, y, Inches(2.2), h, name,
                 size=11, bold=True, color=GREY_DK,
                 anchor=MSO_ANCHOR.MIDDLE)
        track_x = left + Inches(2.3)
        track_w = width - Inches(2.3)
        add_rect(slide, track_x, y + h * 0.3, track_w, h * 0.4, fill=GREY_LT)
        bar_w = Emu(int(track_w * progress))
        if bar_w > Inches(0.3):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, track_x, y + h * 0.2, bar_w, h * 0.6
            )
            fill_shape(shape, color)
        if progress < 0.5:
            mark_x = track_x + bar_w + Inches(0.1)
            warn = slide.shapes.add_shape(
                MSO_SHAPE.LIGHTNING_BOLT, mark_x, y + h * 0.15,
                Inches(0.35), h * 0.7
            )
            fill_shape(warn, ACCENT)


def draw_triple_crown(slide, left, top, width, height, items):
    """Vertical checklist with ✓/!/✗ status icons.

    items = [(label, status), ...] where status in {"ok", "warn", "fail"}.
    Useful for any "audit"-style slide (Triple Crown, capability checklist, etc).
    """
    n = len(items)
    row_h = height / n
    for i, (label, status) in enumerate(items):
        y = top + row_h * i
        ic_w = Inches(0.6)
        ic = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y + Inches(0.05),
                                    ic_w, ic_w)
        if status == "ok":
            fill_shape(ic, GREEN)
            ic_text = "✓"
        elif status == "fail":
            fill_shape(ic, ACCENT)
            ic_text = "✗"
        else:
            fill_shape(ic, AMBER)
            ic_text = "!"
        tf = ic.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = ic_text
        r.font.name = FONT
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = WHITE
        add_text(slide, left + ic_w + Inches(0.25), y,
                 width - ic_w - Inches(0.3), row_h, label,
                 size=14, bold=True, color=GREY_DK,
                 anchor=MSO_ANCHOR.MIDDLE)


def draw_three_columns(slide, left, top, width, height, columns, num_color=NAVY):
    """Three numbered cards side-by-side. Standard layout for playbooks /
    three-step recommendations / three-pillar frameworks.

    columns = [(title, body_lines), ...]
    body_lines is a list passed straight to add_bullets — strings or
    (bold_part, rest) tuples.
    """
    n = len(columns)
    gap = Inches(0.25)
    col_w = (width - gap * (n - 1)) / n

    for i, (title, body) in enumerate(columns):
        x = left + (col_w + gap) * i
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, top,
                                       Inches(0.55), Inches(0.55))
        fill_shape(badge, num_color)
        tf = badge.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.name = FONT
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = WHITE

        add_text(slide, x + Inches(0.7), top + Inches(0.05),
                 col_w - Inches(0.7), Inches(0.5),
                 title, size=14, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)

        body_top = top + Inches(0.75)
        body_h = height - Inches(0.75)
        add_rect(slide, x, body_top, col_w, body_h,
                 fill=BG_LIGHT, line=GREY_LT)
        add_bullets(slide, x + Inches(0.2), body_top + Inches(0.15),
                    col_w - Inches(0.4), body_h - Inches(0.3),
                    body, size=11, line_spacing=1.2)


# =============================================================================
# PDF (SPEAKER SCRIPTS)
# =============================================================================

# Hex equivalents of the ESADE palette for reportlab
PDF_NAVY    = HexColor("#0A2540")
PDF_ACCENT  = HexColor("#C0392B")
PDF_GREY_DK = HexColor("#1F2937")
PDF_GREY_LT = HexColor("#6B7280")


def make_pdf_doc(path, title, author):
    """Return a ready-to-use SimpleDocTemplate with ESADE margins."""
    return SimpleDocTemplate(
        path,
        pagesize=A4,
        leftMargin=2.2 * cm, rightMargin=2.2 * cm,
        topMargin=2.0 * cm, bottomMargin=2.0 * cm,
        title=title,
        author=author,
    )


def pdf_styles():
    """Return a dict of ParagraphStyles tuned for speaker-script PDFs."""
    base = getSampleStyleSheet()
    return {
        "title": ParagraphStyle(
            "title", parent=base["Title"], fontName="Helvetica-Bold",
            fontSize=20, leading=24, textColor=PDF_NAVY, spaceAfter=4,
        ),
        "subtitle": ParagraphStyle(
            "subtitle", parent=base["Normal"], fontName="Helvetica-Oblique",
            fontSize=11, leading=14, textColor=PDF_GREY_LT, spaceAfter=18,
        ),
        "section": ParagraphStyle(
            "section", parent=base["Heading2"], fontName="Helvetica-Bold",
            fontSize=14, leading=18, textColor=PDF_NAVY,
            spaceBefore=14, spaceAfter=4, keepWithNext=True,
        ),
        "meta": ParagraphStyle(
            "meta", parent=base["Normal"], fontName="Helvetica-Oblique",
            fontSize=10, leading=13, textColor=PDF_ACCENT, spaceAfter=8,
            keepWithNext=True,
        ),
        "body": ParagraphStyle(
            "body", parent=base["Normal"], fontName="Helvetica",
            fontSize=11, leading=16, textColor=PDF_GREY_DK,
            alignment=TA_JUSTIFY, spaceAfter=8,
        ),
        "quote": ParagraphStyle(
            "quote", parent=base["Normal"], leftIndent=18, rightIndent=18,
            fontName="Helvetica", fontSize=11, leading=16, textColor=PDF_GREY_DK,
            spaceBefore=6, spaceAfter=10,
        ),
        "qa_q": ParagraphStyle(
            "qa_q", parent=base["Normal"], fontName="Helvetica-Bold",
            fontSize=11, leading=14, textColor=PDF_NAVY,
            spaceBefore=10, spaceAfter=2, keepWithNext=True,
        ),
        "qa_a": ParagraphStyle(
            "qa_a", parent=base["Normal"], fontName="Helvetica",
            fontSize=11, leading=15, textColor=PDF_GREY_DK,
            leftIndent=14, spaceAfter=6,
        ),
        "footnote": ParagraphStyle(
            "foot", parent=base["Normal"], fontName="Helvetica-Oblique",
            fontSize=10, leading=13, textColor=PDF_GREY_LT,
            alignment=TA_LEFT,
        ),
    }
